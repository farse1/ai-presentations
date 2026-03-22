import streamlit as st
from pptx import Presentation
from fpdf import FPDF
import fitz  # PyMuPDF
from langchain_openai import ChatOpenAI
from langchain_community.tools.tavily_search import TavilySearchResults
import json
import os
import re

# Set page config once at the very top
st.set_page_config(page_title="AI Presentation & PDF Creator", page_icon="📊", layout="wide")

# --- 1. API KEY LOGIC (Fixes the Duplicate Error) ---
# Check Streamlit Secrets first
secret_openai = st.secrets.get("OPENAI_API_KEY")
secret_tavily = st.secrets.get("TAVILY_API_KEY")

with st.sidebar:
    st.header("⚙️ Configuration")
    
    # OpenAI Key Logic
    if secret_openai:
        o_api = secret_openai
        st.success("✅ OpenAI Key: Loaded from Secrets")
    else:
        o_api = st.text_input("Enter OpenAI API Key", type="password", key="input_openai")
        st.info("To avoid typing this, add it to Streamlit Secrets.")

    # Tavily Key Logic
    if secret_tavily:
        t_api = secret_tavily
        st.success("✅ Tavily Key: Loaded from Secrets")
    else:
        t_api = st.text_input("Enter Tavily API Key", type="password", key="input_tavily")

    st.divider()
    num_slides = st.slider("Number of Slides", 5, 20, 10, key="slide_slider")

# --- 2. HELPER FUNCTIONS ---

def create_pptx(slides_data):
    prs = Presentation()
    for slide_info in slides_data:
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = slide_info.get("title", "Slide")
        slide.placeholders[1].text = slide_info.get("content", "")
    path = "presentation.pptx"
    prs.save(path)
    return path

def create_pdf(slides_data):
    pdf = FPDF(orientation="landscape", unit="mm", format="A4")
    pdf.set_auto_page_break(auto=True, margin=15)
    for slide_info in slides_data:
        pdf.add_page()
        pdf.set_fill_color(200, 220, 255)
        pdf.rect(0, 0, 297, 30, 'F')
        pdf.set_font("Arial", 'B', 24)
        pdf.set_xy(10, 10)
        pdf.cell(0, 10, slide_info.get("title", "Slide").encode('latin-1', 'replace').decode('latin-1'), ln=True)
        pdf.set_font("Arial", size=14)
        pdf.set_xy(10, 40)
        content = slide_info.get("content", "").encode('latin-1', 'replace').decode('latin-1')
        pdf.multi_cell(0, 10, content)
    path = "presentation.pdf"
    pdf.output(path)
    return path

def extract_text(file):
    try:
        if file.name.endswith("pptx"):
            prs = Presentation(file)
            return "\n".join([shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text")])
        else:
            doc = fitz.open(stream=file.read(), filetype="pdf")
            return "\n".join([page.get_text() for page in doc])
    except Exception:
        return ""

# --- 3. MAIN UI ---
st.title("🚀 AI Presentation & PDF Generator")
st.markdown("Enter a topic and optionally upload a reference. I'll search the web and create your files.")

topic = st.text_input("Topic:", placeholder="e.g., The impact of quantum computing on cybersecurity", key="topic_input")
ref_file = st.file_uploader("Upload reference (Optional)", type=["pdf", "pptx"], key="file_uploader")

if st.button("Generate Presentation", type="primary"):
    if not o_api or not t_api or not topic:
        st.error("❌ Error: Missing API Keys or Topic.")
    else:
        try:
            with st.spinner("🔍 Researching the web and analyzing references..."):
                os.environ["TAVILY_API_KEY"] = t_api
                
                # Context 
                ref_text = extract_text(ref_file)[:1500] if ref_file else ""
                search = TavilySearchResults(max_results=3)
                web_data = search.invoke(topic)
                
                # LLM Generation
                llm = ChatOpenAI(model="gpt-4o-mini", api_key=o_api, temperature=0.7)
                prompt = f"""
                Create a detailed presentation on: {topic}
                Using reference text: {ref_text}
                And web data: {web_data}
                
                Create exactly {num_slides} slides.
                Return ONLY a JSON array of objects with "title" and "content" keys.
                Example: [{{"title": "Title", "content": "Point 1\\nPoint 2"}}]
                """
                
                res = llm.invoke(prompt)
                
                # Clean JSON string
                json_match = re.search(r"\[.*\]", res.content, re.DOTALL)
                if not json_match:
                    raise ValueError("AI did not return valid JSON format.")
                
                slides_json = json.loads(json_match.group())
                
                # Create Files
                pptx_path = create_pptx(slides_json)
                pdf_path = create_pdf(slides_json)
                
                st.success(f"✅ Successfully created {len(slides_json)} slides!")
                
                # Download Buttons
                col1, col2 = st.columns(2)
                with col1:
                    with open(pptx_path, "rb") as f:
                        st.download_button("📥 Download PPTX", f, file_name=f"{topic}.pptx", key="dl_pptx")
                with col2:
                    with open(pdf_path, "rb") as f:
                        st.download_button("📥 Download PDF", f, file_name=f"{topic}.pdf", key="dl_pdf")

        except Exception as e:
            st.error(f"⚠️ An error occurred: {str(e)}")
